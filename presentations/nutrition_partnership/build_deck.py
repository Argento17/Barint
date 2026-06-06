# -*- coding: utf-8 -*-
"""
Bari — Nutrition Expert Partnership Deck
Single source of truth: SLIDES list -> emits both a .pptx and a markdown spec.

All numbers are grounded in real repo artifacts (frontend JSON corpus,
TASK registry, evidence registry, memory). Nothing about product/nutrition
data is invented.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS  = os.path.join(OUT_DIR, "assets")
LOGO_LIGHT = os.path.join(ASSETS, "logo_light.png")   # for dark backgrounds? no -> ink wordmark for light bg
LOGO_DARK  = os.path.join(ASSETS, "logo_dark.png")    # white wordmark for dark bg

def _img(name):
    return os.path.join(ASSETS, name)

def _aspect(path):
    from PIL import Image
    with Image.open(path) as im:
        return im.height / im.width

# ---------------------------------------------------------------- palette
INK      = RGBColor(0x0E, 0x1A, 0x2B)   # near-black navy
NAVY     = RGBColor(0x14, 0x2A, 0x45)   # deep navy
SLATE    = RGBColor(0x3A, 0x4A, 0x5E)   # body slate
MIST     = RGBColor(0x6B, 0x7A, 0x8C)   # secondary grey
TEAL     = RGBColor(0x1F, 0x9E, 0x8A)   # primary accent (Bari green-teal)
TEALDK   = RGBColor(0x15, 0x6B, 0x5E)
GOLD     = RGBColor(0xC9, 0x9A, 0x3F)   # secondary accent
PAPER    = RGBColor(0xF6, 0xF4, 0xEF)   # warm paper
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LINE     = RGBColor(0xD8, 0xDD, 0xE3)
CHIP_A   = RGBColor(0x1F, 0x9E, 0x8A)
CHIP_BG  = RGBColor(0xEC, 0xF3, 0xF1)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

# ---------------------------------------------------------------- helpers
def _solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def _txt(tf, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=6):
    """runs: list of paragraphs; each paragraph = list of (text,size,color,bold,italic)."""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (t, sz, col, bold, ital) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.italic = ital
            r.font.name = "Segoe UI"
    return tf

def add_box(slide, x, y, w, h):
    return slide.shapes.add_textbox(x, y, w, h).text_frame

def add_rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _solid(sp, color)
    sp.shadow.inherit = False
    return sp

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_logo(slide, dark=False):
    """Bari logo, top-right corner — present on every slide (master-style)."""
    path = LOGO_DARK if dark else LOGO_LIGHT
    slide.shapes.add_picture(path, Inches(11.18), Inches(0.30), width=Inches(1.7))

def add_pageno(slide, n, dark=False):
    """Page number, bottom-right, on every slide."""
    col = RGBColor(0x9F, 0xB0, 0xC0) if dark else MIST
    tf = add_box(slide, Inches(12.35), Inches(7.02), Inches(0.8), Inches(0.4))
    _txt(tf, [[(str(n), 11, col, False, False)]], align=PP_ALIGN.RIGHT)

def place_right_image(slide, name, x=Inches(8.0), y=Inches(1.62), w=Inches(4.75)):
    path = _img(name)
    h = Emu(int(w * _aspect(path)))
    # keep within footer; cap height
    maxh = Inches(4.95)
    if h > maxh:
        h = maxh; w = Emu(int(h / _aspect(path))); x = Inches(8.25)
    slide.shapes.add_picture(path, x, y, width=w, height=h)

def place_band_image(slide, name, w=Inches(12.0), bottom=Inches(7.28), xc=True):
    path = _img(name)
    h = Emu(int(w * _aspect(path)))
    maxh = Inches(2.0)                       # keep clear of column text above
    if h > maxh:
        h = maxh; w = Emu(int(h / _aspect(path)))
    y = Emu(int(bottom - h))
    x = Emu(int((EMU_W - w) / 2)) if xc else Inches(0.66)
    slide.shapes.add_picture(path, x, y, width=w, height=h)

# ---------------------------------------------------------------- layouts
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def cover_slide(prs, s):
    sl = blank(prs)
    add_rect(sl, 0, 0, EMU_W, EMU_H, INK)
    add_rect(sl, 0, 0, Inches(0.28), EMU_H, TEAL)
    add_rect(sl, Inches(0.9), Inches(2.55), Inches(2.2), Inches(0.06), TEAL)
    tf = add_box(sl, Inches(0.9), Inches(1.5), Inches(11.4), Inches(1.0))
    _txt(tf, [[("BARI", 22, TEAL, True, False),
               ("  ·  Nutrition Intelligence", 22, MIST, False, False)]])
    tf = add_box(sl, Inches(0.85), Inches(2.75), Inches(11.6), Inches(2.2))
    _txt(tf, [[(s["title"], 46, WHITE, True, False)],
              [(s["subtitle"], 22, RGBColor(0xB8,0xC4,0xD0), False, True)]],
         space_after=14)
    tf = add_box(sl, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.7))
    _txt(tf, [[(s["footer"], 13, MIST, False, False)]])
    add_logo(sl, dark=True)
    add_notes(sl, s["notes"])
    return sl

def section_slide(prs, s):
    sl = blank(prs)
    add_rect(sl, 0, 0, EMU_W, EMU_H, NAVY)
    add_rect(sl, 0, Inches(2.9), EMU_W, Inches(0.04), TEAL)
    tf = add_box(sl, Inches(0.9), Inches(1.65), Inches(11.5), Inches(0.6))
    _txt(tf, [[("SECTION %s" % s["num"], 18, TEAL, True, False)]])
    tf = add_box(sl, Inches(0.85), Inches(3.15), Inches(11.6), Inches(2.4))
    _txt(tf, [[(s["title"], 40, WHITE, True, False)],
              [(s["sub"], 19, RGBColor(0xB8,0xC4,0xD0), False, True)]], space_after=14)
    add_logo(sl, dark=True)
    add_notes(sl, s["notes"])
    return sl

def content_slide(prs, s):
    sl = blank(prs)
    add_rect(sl, 0, 0, EMU_W, EMU_H, WHITE)
    add_rect(sl, 0, 0, EMU_W, Inches(1.25), PAPER)
    add_rect(sl, 0, Inches(1.25), EMU_W, Inches(0.03), LINE)
    add_rect(sl, Inches(0.55), Inches(0.34), Inches(0.12), Inches(0.62), TEAL)
    # kicker + title
    tf = add_box(sl, Inches(0.85), Inches(0.26), Inches(11.6), Inches(0.45))
    _txt(tf, [[(s.get("kicker","").upper(), 12, TEALDK, True, False)]])
    tf = add_box(sl, Inches(0.83), Inches(0.6), Inches(10.9), Inches(0.7))
    _txt(tf, [[(s["title"], 27, INK, True, False)]])

    img = s.get("img"); mode = s.get("img_mode")
    right = img and mode == "right"
    band  = img and mode == "band"
    text_w = Inches(6.75) if right else Inches(11.6)

    body_top = 1.55
    # optional lead line
    if s.get("lead"):
        tf = add_box(sl, Inches(0.85), Inches(body_top), text_w, Inches(0.75))
        _txt(tf, [[(s["lead"], 15, SLATE, False, True)]])
        body_top += 0.66

    bullets = s.get("bullets", [])
    cols = s.get("columns")
    if cols:
        _two_col(sl, cols, body_top, compact=band)
    elif s.get("stats"):
        _stat_band(sl, s["stats"], body_top)
        if bullets:
            _bullet_block(sl, bullets, Inches(0.85), Inches(body_top+1.5), text_w)
    else:
        _bullet_block(sl, bullets, Inches(0.85), Inches(body_top), text_w)

    # embedded visual
    if right:
        place_right_image(sl, img)
    elif band:
        place_band_image(sl, img)
    elif s.get("visual"):
        # fallback: textual visual recommendation strip
        add_rect(sl, 0, Inches(6.7), EMU_W, Inches(0.8), PAPER)
        add_rect(sl, 0, Inches(6.7), EMU_W, Inches(0.03), LINE)
        tf = add_box(sl, Inches(0.85), Inches(6.82), Inches(11.6), Inches(0.6))
        _txt(tf, [[("VISUAL  ", 11, TEALDK, True, False),
                   (s["visual"], 12, SLATE, False, True)]])
    add_logo(sl, dark=False)
    add_notes(sl, s["notes"])
    return sl

def _bullet_block(sl, bullets, x, y, w):
    tf = add_box(sl, x, y, w, Inches(5.0))
    paras = []
    for b in bullets:
        if isinstance(b, tuple):
            head, sub = b
            paras.append([("▸  ", 15, TEAL, True, False),
                          (head+"  ", 15, INK, True, False),
                          (sub, 14, SLATE, False, False)])
        else:
            paras.append([("▸  ", 15, TEAL, True, False),
                          (b, 15, INK, False, False)])
    _txt(tf, paras, space_after=11)

def _two_col(sl, cols, body_top, compact=False):
    xs = [Inches(0.85), Inches(6.95)]
    w = Inches(5.5)
    body_h = Inches(2.5) if compact else Inches(4.4)
    fs_h, fs_b, fs_s = (14, 12.5, 11.5) if compact else (15, 13.5, 12.5)
    sp = 5 if compact else 8
    for i, col in enumerate(cols):
        x = xs[i]
        if col["h"]:
            add_rect(sl, x, Inches(body_top), w, Inches(0.5), CHIP_BG)
            add_rect(sl, x, Inches(body_top), Inches(0.1), Inches(0.5), TEAL)
            tf = add_box(sl, x+Inches(0.2), Inches(body_top+0.06), w-Inches(0.3), Inches(0.4))
            _txt(tf, [[(col["h"], fs_h, TEALDK, True, False)]])
            items_y = body_top + 0.66
        else:
            items_y = body_top
        tf = add_box(sl, x+Inches(0.05), Inches(items_y), w-Inches(0.1), body_h)
        paras = []
        for b in col["items"]:
            if isinstance(b, tuple):
                paras.append([("•  ", fs_b, TEAL, True, False),
                              (b[0]+"  ", fs_b, INK, True, False),
                              (b[1], fs_s, SLATE, False, False)])
            else:
                paras.append([("•  ", fs_b, TEAL, True, False),
                              (b, fs_b, INK, False, False)])
        _txt(tf, paras, space_after=sp)

def _stat_band(sl, stats, body_top):
    n = len(stats)
    gap = Inches(0.25)
    total_w = Inches(11.6)
    cw = Emu(int((total_w - gap*(n-1)) / n))
    x = Inches(0.85)
    for st in stats:
        add_rect(sl, x, Inches(body_top), cw, Inches(1.3), PAPER)
        add_rect(sl, x, Inches(body_top), cw, Inches(0.08), TEAL)
        tf = add_box(sl, x+Inches(0.1), Inches(body_top+0.18), cw-Inches(0.2), Inches(0.7))
        _txt(tf, [[(st["big"], 30, TEALDK, True, False)]], align=PP_ALIGN.CENTER)
        tf = add_box(sl, x+Inches(0.1), Inches(body_top+0.85), cw-Inches(0.2), Inches(0.45))
        _txt(tf, [[(st["cap"], 11.5, SLATE, False, False)]], align=PP_ALIGN.CENTER)
        x = Emu(x + cw + gap)

# ================================================================ CONTENT
SLIDES = []
def S(**k): SLIDES.append(k)

# ---- COVER
S(kind="cover",
  title="Bari — Building Nutrition Intelligence",
  subtitle="A partnership conversation for a senior nutrition scientist",
  footer="Confidential  ·  Prepared for nutrition-expert evaluation  ·  June 2026",
  notes=("Open warm but get to the point in 30 seconds. 'You know nutrition deeply. "
         "You've also watched the public information layer around food rot — scores that "
         "contradict each other, labels nobody can read, influencers filling the vacuum. "
         "Bari is an attempt to fix that layer with real engineering and real evidence "
         "discipline. I'm not here to pitch you a logo — I'm here to show you a system that "
         "already exists, and to ask whether you want to put your name and judgement inside it "
         "while it's still being shaped.' Set expectation: ~25 slides, heavy on what's built, "
         "honest about what isn't."))

# ============================== SECTION 1 — PROBLEM
S(kind="section", num="1", title="The Problem",
  sub="The consumer nutrition-information layer is broken — and trust is leaving with it.",
  notes=("Frame the problem as a SCIENTIST would feel it, not a marketer. The expert has "
         "personally seen patients act on garbage scores. Validate that frustration first; "
         "it earns the right to everything after."))

S(kind="content", kicker="The Problem · 1 of 5",
  title="Nutrition information is structurally broken",
  lead="The failure is not a lack of data. It is the absence of a trustworthy interpretation layer between the label and the person.",
  bullets=[
    ("Data exists, meaning doesn't.", "Every package has a panel; almost no consumer can convert it into a decision."),
    ("Interpretation is outsourced to whoever shouts loudest.", "Influencers and brands fill the vacuum the science left open."),
    ("Context collapses.", "A number with no category, no confidence, and no reasoning is noise dressed as authority."),
  ],
  visual="Split image — a dense ingredient panel on the left, a confused shopper on the right; one arrow between them labelled 'missing layer'.",
  notes=("Land the core thesis early: the gap is interpretation, not raw data. Everything Bari "
         "builds sits in that gap. Don't rush — let the scientist nod at 'meaning doesn't exist'."))

S(kind="content", kicker="The Problem · 2 of 5",
  title="Consumers are confused — predictably so",
  bullets=[
    ("Marketing language mimics health language.", "‘טבעי’, ‘ללא’, ‘עשיר ב־’ carry no consistent meaning."),
    ("Front-of-pack claims contradict the back-of-pack reality.", "The diet-yogurt that scores worse than the full-fat one."),
    ("Health-by-omission.", "‘No sugar’ says nothing about what replaced it."),
    ("No category sense.", "Shoppers compare a cracker to a yogurt because nothing tells them not to."),
  ],
  visual="Three real front-of-pack claims with a red 'says nothing about…' annotation under each.",
  notes=("Use a concrete Israeli-shelf example the expert recognises: a 'diet' dessert that is "
         "engineered worse than its plain sibling. This previews our maadanim (מילקי = E) finding "
         "without naming the section yet."))

S(kind="content", kicker="The Problem · 3 of 5",
  title="Ingredient lists are effectively unreadable",
  bullets=[
    ("Length is weaponised.", "20+ ingredients, additives behind code numbers, function hidden."),
    ("Order obscures dose.", "Position implies quantity but never states it."),
    ("Named-additive identity is lost.", "‘חלבון חלב’, emulsifiers, stabilisers read as one grey blur."),
    "Even a trained dietitian needs minutes per product. A shopper has seconds.",
  ],
  visual="A real long ingredient list with additive codes highlighted and a stopwatch overlay ('avg shopper attention: 6s').",
  notes=("This is where you signal technical seriousness: Bari's engine actually parses named "
         "additives and reasons about them (named-additive identity, engine 0.4.0). Tease, don't "
         "explain yet."))

S(kind="content", kicker="The Problem · 4 of 5",
  title="Existing scores fail — and fail quietly",
  columns=[
    {"h":"Why current scores break", "items":[
        ("One global formula.", "A bread rule applied to yogurt is wrong by construction."),
        ("No confidence.", "A guess and a verified result look identical."),
        ("No reasoning.", "A letter with no ‘why’ can't be trusted or challenged."),
        ("Gamed.", "Reformulate to beat the metric, not to improve the food."),
    ]},
    {"h":"What a credible score needs", "items":[
        ("Category-relative judgement.", "Compare like with like."),
        ("Explicit confidence.", "Verified / partial / insufficient — never hidden."),
        ("Traceable reasoning.", "Every grade defends itself from real signals."),
        ("Evidence-anchored rules.", "Not opinion, not vibes."),
    ]},
  ],
  visual="Two-column 'Today vs. What's required' table; left column greyed, right column in Bari teal.",
  notes=("This slide quietly defines the evaluation criteria the expert will use on US — and we "
         "happen to meet all four. Let them feel they wrote the requirements."))

S(kind="content", kicker="The Problem · 5 of 5",
  title="Trust is deteriorating — that's the opening",
  lead="When institutions go quiet, the loudest voice wins. The result is a market actively hungry for a source that is rigorous and readable at the same time.",
  bullets=[
    ("The authority vacuum is real.", "Regulators are slow; brands are conflicted; social media is fast and wrong."),
    ("Consumers know they're being sold to.", "Skepticism is high — credibility is scarce and therefore valuable."),
    ("The winner is whoever is both trusted AND usable.", "Rigor alone is ignored; usability alone is dangerous. Bari targets the intersection."),
  ],
  visual="Trust-vs-time line trending down for ‘brands/influencers’, with a gap labelled ‘the position Bari is building toward’.",
  notes=("Transition line: 'So the question isn't whether this layer is needed — it's whether "
         "anyone can build it credibly. Here's what we are building.'"))

# ============================== SECTION 2 — VISION
S(kind="section", num="2", title="The Bari Vision",
  sub="One rigorous engine powering three audiences — consumer, professional, industry — and a trusted commerce marketplace.",
  notes=("Shift energy from problem to ambition. Keep it grounded: vision slides are where "
         "skeptics disengage if you over-claim. Anchor every layer to something that already exists."))

S(kind="content", kicker="The Vision",
  title="What Bari is",
  lead="Bari is a nutrition-intelligence platform: an evidence-anchored engine that turns messy real-world food labels into honest, category-aware, explainable judgement — and the trusted place to act on it.",
  bullets=[
    ("Not a blog. Not a single score.", "An interpretation engine with governance, confidence and provenance built in."),
    ("Starts consumer-facing in Hebrew.", "Live comparison pages on the Israeli shelf today."),
    ("Three intelligence audiences, one engine.", "Consumer, professional and industry — same engine, progressively deeper surfaces."),
    ("And a trusted marketplace.", "Bari is also where trust converts to action — discovering and buying better products — without ever putting a grade up for sale."),
  ],
  visual="One-engine / four-surface diagram: central engine (BSIP) feeding Consumer, Professional, Industry and a Commerce/Marketplace layer.",
  notes=("Repeat the spine: ONE engine, THREE audiences, PLUS a marketplace. The consumer site is "
         "the proof and the data flywheel; professional + industry are where defensibility lives; "
         "the marketplace is where trust becomes revenue. Stress the integrity line: we monetise "
         "the transaction, never the grade."))

S(kind="content", kicker="The Vision",
  title="Three audiences + a marketplace, one engine",
  lead="The same engine serves three intelligence audiences — and a commerce layer that turns earned trust into transactions.",
  columns=[
    {"h":"Consumer layer  (live)", "items":[
        "Category comparison pages on the real Israeli shelf",
        "Honest grades + plain-language verdicts",
        "Confidence shown, never hidden",
    ]},
    {"h":"Professional + commerce  (next)", "items":[
        "Decision-support for dietitians & clinicians",
        "Category dossiers, evidence trails, defensible reasoning",
        "A marketplace: discover → buy better products (grade stays unsellable)",
    ]},
  ],
  visual="Four horizontal bands (Consumer / Professional / Industry-Data / Commerce-Marketplace) sharing one engine column on the left.",
  notes=("Note explicitly: the DATA layer underneath (proprietary scored corpus + evidence "
         "registry) is what makes the upper layers possible and hard to copy; the marketplace is "
         "what makes the trust pay. The integrity guardrail is structural — you can buy a product "
         "through Bari, you can never buy a better grade."))

S(kind="content", kicker="The Vision · 5-year arc",
  title="A credible five-year path",
  lead="Audience → trust → tools → industry → geography. Each year funds and de-risks the next; no leap is unbacked.",
  bullets=[
    ("Years 1–2 — earn the shelf.", "A live Israeli comparison platform, then deepening category intelligence."),
    ("Year 3 — professional tools.", "Subscriptions for dietitians & clinics, once the engine is category-proven."),
    ("Years 4–5 — commerce, industry & geography.", "Retail integration and a product marketplace, then US expansion + a nutrition-intelligence ecosystem others build on."),
  ],
  visual="Horizontal 5-year roadmap ribbon with a widening 'defensibility / revenue' wedge beneath it.",
  notes=("Stress the LOGIC of the order, not the dates. Audience first because trust is the scarce "
         "asset; professional tools only after the engine has been validated on real categories; "
         "US last because the engine must be proven portable before geography multiplies cost."))

# ============================== SECTION 3 — WHAT WE BUILT
S(kind="section", num="3", title="What We've Already Built",
  sub="This is not a concept deck. A working pipeline, a scoring engine, governance, and a live site already exist.",
  notes=("This is the credibility core. Pace slows here. The expert must leave this section "
         "thinking 'this is further along and more disciplined than I expected.'"))

S(kind="content", kicker="What We Built · pipeline",
  title="A three-stage data pipeline: BSIP0 → BSIP1 → BSIP2",
  columns=[
    {"h":"Ingestion & consolidation", "items":[
        ("BSIP0 — acquisition.", "Real retailer scraping (Shufersal), label & nutrition parsing, raw-source persistence for offline replay."),
        ("BSIP1 — trust layer.", "Cross-retailer consolidation, observation-quality + canonical trust scoring, 10-field semantic enrichment. Not scoring."),
    ]},
    {"h":"Scoring & interpretation", "items":[
        ("BSIP2 — the engine.", "Four interpretive layers: Structural, Nutritional, Metabolic, Engineering."),
        ("Universal core + archetypes.", "One core, category-specific interpretation via a 3-stage router (anchor → signal → resolution). Engine v0.4.0."),
    ]},
  ],
  visual="Left-to-right pipeline diagram with three labelled stages and a 'raw source preserved' callout under BSIP0.",
  notes=("Key message: scoring is DELIBERATELY isolated downstream. Ingestion, consolidation, and "
         "judgement are separate layers — so the expert can audit the judgement layer without "
         "touching plumbing. That separation is itself a rigor signal."))

S(kind="content", kicker="What We Built · rigor infrastructure",
  title="Evidence, confidence and validation are built in — not bolted on",
  bullets=[
    ("Evidence Registry.", "20 food findings + 20 guardrails + 20 deferred domains; every rule traces to a catalogued finding, not opinion. Each finding has a tracking id so we can say exactly why a rule exists."),
    ("Confidence system.", "Three explicit states — verified / partial / insufficient. When the label can't support a score we publish no score at all, rather than fake certainty."),
    ("Validation system.", "A fixed set of 12 hand-checked reference products plus an automated regression test that re-scores them on every change — so a parsing mistake is caught before it can reach the site."),
    ("Governance.", "A written constitution: 6-article comparison governance, use-case guardrails, an exception registry for every deliberate deviation."),
  ],
  visual="Four-quadrant 'rigor stack': Evidence / Confidence / Validation / Governance, each with its real artifact name.",
  notes=("This is the slide that separates Bari from every content startup. Emphasise: we have a "
         "MECHANISM to refuse to answer (null on insufficient) and a MECHANISM to record every "
         "exception. A scientist trusts a system that knows its own limits."))

S(kind="content", kicker="What We Built · surfaces & operations",
  title="A live platform and an operating model behind it",
  columns=[
    {"h":"Product surfaces (live)", "items":[
        ("Category framework.", "Repeatable factory: shelf-map → corpus → score → QA → publish."),
        ("Comparison platform.", "Live Hebrew RTL pages with grades, verdicts, expandable reasoning."),
        ("External data layer.", "10 read-only authoritative clients — OFF, gov data, PubMed, DSLD, Tzameret…"),
    ]},
    {"h":"Operating model", "items":[
        ("Command Center.", "Authoritative task registry → derived dashboard; nothing ships untracked."),
        ("Agent OS.", "Specialist agents (nutrition, data, QA, content) with a verification gate before close."),
        ("Provenance discipline.", "Frozen invariants, dated rulings, no silent score changes."),
    ]},
  ],
  visual="Screenshot of a live comparison page beside a Command Center dashboard view.",
  notes=("The operating model is the unsung moat: an auditable trail from a raw label to a "
         "published grade, with a registry that records who decided what and why. Most teams have "
         "none of this. Mention the 10 live external clients as evidence we reason from "
         "authoritative sources, not scraping guesses."))

S(kind="content", kicker="What We Built · the moat",
  title="Why this is hard to replicate",
  bullets=[
    ("It's not the model — it's the accumulated judgement.", "Every category adds calibrated rules, evidence links and stress-tested edge cases that don't transfer for free."),
    ("Real-shelf data is messy and earned.", "Parsing real Israeli labels (and surviving their traps) took real failures to fix — e.g. the fat-overwrite bug below."),
    ("Governance compounds.", "An exception registry and frozen invariants mean the system gets MORE trustworthy over time, not noisier."),
    ("Honesty is a feature competitors can't fake late.", "‘We score this null because we can't verify it’ is a trust posture you build from day one."),
  ],
  visual="Layered-moat diagram: Data → Calibration → Evidence → Governance → Trust, each ring wider and harder to cross.",
  notes=("Tie replication difficulty to TIME and DISCIPLINE, not secrecy. Anyone can copy a UI in "
         "a week; nobody can copy two years of calibrated category judgement and a governance "
         "trail. This is the slide an investor remembers."))

# ============================== SECTION 4 — SCIENTIFIC DEPTH
S(kind="section", num="4", title="Scientific Depth",
  sub="The methodology a nutrition scientist would actually respect — shown through real rulings, not slogans.",
  notes=("Now speak peer-to-peer. This section is for the scientist's internal credibility test. "
         "Lead with philosophy, then prove each principle with a real category ruling."))

S(kind="content", kicker="Scientific Depth · principles",
  title="The methodology, stated plainly",
  bullets=[
    ("Evidence-based.", "Rules derive from a registry of findings, each with an id and rationale — challengeable, versioned."),
    ("Label-observability.", "We only score what the label actually lets us observe. No invented values; unverifiable → insufficient."),
    ("Food reality over opinion.", "We interpret the food as engineered, not the marketing story around it."),
    ("Processing read with nuance.", "NOVA-aware but not NOVA-blind; named additives and matrix effects reasoned about, not just counted."),
    ("Category-specific calibration.", "Protein, fat, fortification all judged category-relative — a yogurt rule never touches bread."),
    ("Stress-tested.", "Every category passes an adversarial governance audit before it can go live."),
  ],
  visual="Six-principle hexagon, each tile linking to the real example slide that proves it.",
  notes=("Read these as a manifesto, then say: 'Principles are cheap — let me show you five times "
         "this changed an actual grade.' Hand over to the examples."))

S(kind="content", kicker="Scientific Depth · real ruling 1",
  title="Hummus — defining the prepared/raw boundary correctly",
  lead="A naïve engine splits hummus by protein or by the word ‘סלט’. Both are wrong.",
  bullets=[
    ("The real boundary is tahini + sodium + energy density.", "Not protein, never the label noun. Tahini-based salads are correctly kept in scope."),
    ("Result: a defensible 64-product comparison.", "Recalibration moved the distribution toward truth (B-tier widened as cultures/quality were credited correctly)."),
    ("Why it matters scientifically.", "It shows the engine reasons about food structure, not keywords."),
  ],
  visual="Scatter of hummus products on tahini×sodium axes with the prepared/raw boundary drawn through it.",
  notes=("This is the 'they actually think about food' proof. The protein-trap correction "
         "(raw-vs-prepared boundary) is exactly the kind of subtle call a dietitian respects and a "
         "keyword classifier gets wrong."))

S(kind="content", kicker="Scientific Depth · real ruling 2",
  title="Yogurts — crediting live cultures, and refusing to overclaim",
  bullets=[
    ("We found a real bug in our own pipeline.", "Live cultures were being detected when we read the label, but the scoring step looked at a different list and never counted them — 49 products had cultures detected, 0 credited."),
    ("We fixed it at the source.", "We made the scorer read the same culture list; 12 products correctly moved from C up to B once their live cultures actually counted."),
    ("Then we held the line.", "Even after the fix, no yogurt reached an A. A score of 78.7 (a high B) is the honest ceiling — the cleanest candidates simply don't publish enough on their label to justify more."),
  ],
  visual="Before/after grade-distribution bars for yogurt with the two corrected steps annotated.",
  notes=("This is the most persuasive slide for a skeptic: we improved scores when evidence "
         "justified it AND refused to inflate to an A when the data didn't support one. That "
         "asymmetry — quick to credit, slow to crown — is the whole trust thesis in one example."))

S(kind="content", kicker="Scientific Depth · real ruling 3",
  title="Cereals — testing assumptions instead of inheriting them",
  bullets=[
    ("Assumption: fortification is endemic and should be discounted.", "We measured it. Only 27.2% — not endemic. The assumption was wrong, so we didn't apply the discount."),
    ("Router honesty.", "We declared a live NO-GO on cereals because of a cereal-anchor routing gap (misroute rate too high) — rather than shipping a shaky category."),
    ("Discipline over coverage.", "We would rather withhold a category than publish one we can't defend."),
  ],
  visual="Bar showing 'assumed fortification' vs 'measured 27.2%', plus a red NO-GO stamp on the cereals tile.",
  notes=("Use this to prove intellectual honesty about our OWN priors. We let data overturn a "
         "convenient assumption, and we benched a category rather than ship misrouted scores. "
         "That's the behaviour the expert wants to see governing their name."))

S(kind="content", kicker="Scientific Depth · real rulings 4 & 5",
  title="Dairy & fermentation — calibration with frozen invariants",
  columns=[
    {"h":"Dairy ruling", "items":[
        ("Milk = run_004, frozen.", "Top tier 85/A is whole / 4% / goat dairy — and that ceiling does not silently move."),
        ("Category thresholds are explicit.", "Protein, light-claim and calcium thresholds defined, documented, defensible."),
        ("Diet ≠ better.", "Engineered ‘light’ products don't get a free pass."),
    ]},
    {"h":"Fermentation ruling", "items":[
        ("Bread fermentation read with care.", "Sourdough/long-ferment signals credited where observable, not assumed from a word."),
        ("Cheese cultures credited like yogurt.", "Consistent fermentation logic across categories."),
        ("The fat-overwrite catch (EV-029).", "A parser bug zeroed total fat from a trans sub-row across all nutritionList categories — found, fixed centrally, QA-guarded, re-scored."),
    ]},
  ],
  visual="Two mini-cards: a frozen 'milk 85/A' chip with a lock icon; a bug→fix→guard timeline for EV-029.",
  notes=("EV-029 is your strongest engineering-credibility anecdote: a subtle Hebrew final-letter "
         "parsing trap silently corrupted fat values; we caught it, fixed it once at the shared "
         "layer, added a permanent QA guard (COV-006), and re-scored affected categories. That's "
         "how a serious lab behaves. Frozen invariants show we don't let scores drift to flatter "
         "ourselves."))

# ============================== SECTION 5 — PROGRESS
S(kind="section", num="5", title="Current Progress",
  sub="Concrete output to date — categories live, products scored, lessons banked.",
  notes=("Quantify. The skeptic is asking 'is this real or a demo?'. Numbers from the live corpus "
         "answer that."))

S(kind="content", kicker="Current Progress · by the numbers",
  title="What's live today",
  lead="6 live categories · 250+ curated scored products on the real shelf · 10 authoritative external data clients (live-verified).",
  bullets=[
    ("Per category, end-to-end.", "bread 24 · cheese ~52 · hummus 64 · dairy-desserts 84 · bars 18 · yogurt 11 — each scraped, scored, QA'd, published."),
    ("Funnels are honest.", "bread: 256 scanned → 81 scored → 31 curated; maadanim: 200 scraped → 169 scored. We publish the defensible subset, not everything."),
    ("Next engine in flight.", "A supplement-intelligence engine (SIE) — sibling to BSIP2 — is in methodology phase."),
  ],
  visual="Six category tiles each showing n-products + grade-distribution mini-bar; a totals strip beneath.",
  notes=("Be precise and modest: these are CURATED counts (the publishable subset), and you should "
         "say so — it reinforces that we withhold what we can't defend. The 256→81→31 bread funnel "
         "is the single best honesty proof point; lead with it if pressed."))

S(kind="content", kicker="Current Progress · milestones",
  title="Milestones and lessons banked",
  columns=[
    {"h":"Milestones", "items":[
        "Repeatable category factory proven across 6 shelves",
        "Real-retailer ingestion with raw-source replay",
        "Evidence registry + governance constitution in force",
        "Command Center + agent operating model running daily",
        "10-client external evidence layer live-verified (Jun 2026)",
    ]},
    {"h":"Lessons learned", "items":[
        ("Parsing reality is the hard part.", "Real labels hide traps — e.g. a Hebrew-text quirk once let a trans-fat sub-line silently overwrite total fat. Found, fixed once centrally, permanently guarded."),
        ("Detecting a signal isn't crediting it.", "A value read at one step but ignored at the next is invisible — pipelines must be wired end-to-end (the yogurt-culture fix)."),
        ("Measure your priors.", "We assumed cereal fortification was everywhere; it was only 27% — so we dropped the penalty."),
        ("Withholding is a feature.", "Benching a category we can't defend yet beats shipping a shaky one."),
    ]},
  ],
  visual="A horizontal timeline (categories shipped over time) with lesson callouts pinned to the runs that taught them.",
  notes=("Lessons matter more than milestones to a scientist — they prove we learn under contact "
         "with reality. Each lesson maps to a real EV-id or run, so none of this is abstract."))

# ============================== SECTION 6 — WHY WIN
S(kind="section", num="6", title="Why Bari Can Win",
  sub="Defensibility comes from compounding assets, not a single trick.",
  notes=("Switch to strategy. The expert is also evaluating whether this is worth their reputation. "
         "Defensibility is the answer."))

S(kind="content", kicker="Why Bari Can Win",
  title="Six compounding moats",
  columns=[
    {"h":"", "items":[
        ("Proprietary datasets.", "Scored real-shelf corpus that took real failures to clean."),
        ("Category intelligence.", "Calibrated, stress-tested rules per category — not transferable for free."),
        ("Evidence infrastructure.", "A registry that links every rule to a finding."),
    ]},
    {"h":"", "items":[
        ("Calibration engine.", "The judgement layer improves with every category and edge case."),
        ("Consumer trust.", "Built by showing confidence and withholding — hard to retrofit."),
        ("Content flywheel.", "Each category page draws audience → feedback → better calibration → more pages."),
    ]},
  ],
  visual="Flywheel diagram: Audience → Data → Calibration → Trust → more Audience, with the six moats labelled on the rim.",
  notes=("The punchline: these reinforce each other. Trust drives audience; audience drives data; "
         "data sharpens calibration; calibration deepens trust. A competitor must out-run a "
         "spinning flywheel, not a standing start. The expert STRENGTHENS the two softest spokes "
         "— evidence infrastructure and trust."))

# ============================== SECTION 7 — MONETIZATION
S(kind="section", num="7", title="Monetization",
  sub="A staged path from audience to marketplace to ecosystem — each stage earns the right to the next.",
  notes=("Be sober here. A scientist distrusts a revenue fairy tale. Pair every stage with its "
         "real risk. Credibility now buys believability later."))

S(kind="content", kicker="Monetization · the staircase",
  title="Six stages, each de-risking the next",
  lead="Guardrail on every stage: we monetise the transaction, never the grade.",
  columns=[
    {"h":"Stages 1–3", "items":[
        ("1 · Audience + trust.", "Free comparisons. Cheap reach; slow, indirect revenue."),
        ("2 · Content + influence.", "Editorial authority. Compounds trust; mustn't sell out."),
        ("3 · Professional subs.", "Tools for clinics. Real willingness-to-pay; needs depth."),
    ]},
    {"h":"Stages 4–6", "items":[
        ("4 · Commerce / marketplace.", "Discover → buy better products. Scales with trust; stays grade-blind."),
        ("5 · Industry intelligence.", "Shelf insight for retailers. High-value; conflicts governed hard."),
        ("6 · Retail / ecosystem.", "Engine licensing. Durable, defensible; execution + geography cost."),
    ]},
  ],
  visual="A six-step staircase rising left-to-right (marketplace = step 4, highlighted); a rising revenue/defensibility trend above.",
  notes=("Hit the integrity guardrail hard: across EVERY stage — especially the marketplace and "
         "industry stages — we monetise the transaction, never the grade. The instant a brand can "
         "buy a better score, consumer trust (and the expert's reputation) evaporates. Our "
         "governance makes that boundary structural, not a promise. The marketplace is how trust "
         "becomes revenue without compromising the science."))

# ============================== SECTION 8 — ROLE OF EXPERT
S(kind="section", num="8", title="The Role of the Nutrition Expert",
  sub="The single most important section — where your judgement plugs into the system.",
  notes=("Slow down completely. This is the section the whole deck exists to earn. Make it about "
         "THEM: leverage, not labour. Be explicit about what is and isn't expected."))

S(kind="content", kicker="The Expert's Role · what Bari needs",
  title="What we need from nutrition expertise",
  bullets=[
    ("Methodology review.", "Pressure-test the scoring philosophy and category rules against the literature."),
    ("Evidence evaluation.", "Help grow and grade the evidence registry — what's strong, what's deferred, what's overstated."),
    ("Category calibration.", "Set and defend category-relative thresholds (the dairy/yogurt/hummus kind of calls)."),
    ("Dossier creation.", "Co-author the professional-layer category dossiers that clinicians will rely on."),
    ("Scientific governance.", "Sit over the constitution: approve rulings, adjudicate exceptions, guard the invariants."),
    ("Public credibility.", "Be the named scientific authority that makes ‘trustworthy’ literally true."),
  ],
  visual="Six role-cards arranged around the central engine, each plugging into a real subsystem (registry, router, governance…).",
  notes=("Each responsibility maps to a system that ALREADY EXISTS — so the expert is steering a "
         "running engine, not building from zero. That's a far more attractive ask. Emphasise that "
         "their judgement becomes encoded, versioned and reused — leverage, not one-off consulting."))

S(kind="content", kicker="The Expert's Role · scope honesty",
  title="What is NOT expected",
  columns=[
    {"h":"Not your job", "items":[
        "Not building software or pipelines",
        "Not scraping or cleaning data",
        "Not full-time, day-one (engagement can scale in)",
        "Not rubber-stamping — disagreement is the point",
        "Not personal liability for the platform",
    ]},
    {"h":"How it actually fits", "items":[
        ("Judgement in, encoded once.", "You decide a calibration; the engine applies it consistently forever."),
        ("Async + reviewed.", "Agents prepare evidence packets; you adjudicate, you don't fetch."),
        ("Governance-gated.", "Your sign-off (e.g. D7 co-sign) is a real gate in our process, not decoration."),
    ]},
  ],
  visual="Two columns: a greyed 'not expected' list vs. a teal 'your judgement, amplified' workflow arrow.",
  notes=("Defuse the two fears every expert has: 'will this eat my life?' and 'will my name cover "
         "someone else's shortcuts?'. Answer both directly. Mention that sign-off gates (Product/"
         "expert co-sign) are already wired into the operating model — their authority is real."))

S(kind="content", kicker="The Expert's Role · integration",
  title="How nutrition expertise plugs into the operating model",
  lead="The agent OS already runs a delegate → evidence → review → governed-close loop. The expert sits at the review-and-govern node.",
  bullets=[
    ("Specialist agents do the legwork.", "Data/QA/nutrition agents assemble evidence packets and proposed rulings."),
    ("The expert adjudicates.", "Accept / refine / reject — with reasoning captured in the registry."),
    ("Decisions become durable.", "Every ruling is versioned, dated, traceable; nothing relies on memory."),
    ("Authority is structural.", "Roadmap-impacting items can't close without the required sign-off — the system enforces it."),
  ],
  visual="The delegate→evidence→review→close loop with the expert highlighted at the 'review/govern' stage.",
  notes=("This is reassuring AND flattering: we built the governance gates BEFORE the expert "
         "arrived, so their authority is load-bearing from day one, not aspirational. Show the loop "
         "as already turning."))

# ============================== SECTION 9 — 90 DAYS
S(kind="section", num="9", title="Immediate Opportunities — first 90 days",
  sub="Concrete, scoped projects where the expert moves the needle now.",
  notes=("Make it tangible and near. A skeptic converts when they can see week-one work. Offer a "
         "menu, not a mandate."))

S(kind="content", kicker="First 90 Days · concrete projects",
  title="Five places you'd make an immediate difference",
  bullets=[
    ("Supplement Intelligence Engine (SIE).", "In methodology phase now — co-define the (active, dose, form, evidence) unit and the 5 dimensions (Evidence/Dose/Form/Honesty/Safety). Shape the ‘No-Necessity Rule’ before it sets."),
    ("Dairy calibration.", "Ratify or revise the milk/yogurt/cheese thresholds and the frozen invariants. High-trust, high-visibility."),
    ("Evidence registry expansion.", "Grade and extend beyond the current 20 findings — your literature fluency is the rate-limiter here."),
    ("Category reviews.", "Adjudicate live categories (cereals NO-GO, hummus boundary) and clear the next ones to ship."),
    ("Methodology validation.", "Stress-test the whole scoring philosophy and co-sign it — turning ‘rigorous’ into a claim we can defend publicly."),
  ],
  visual="A 90-day board: 5 swimlanes (SIE / Dairy / Evidence / Reviews / Methodology) with week-1 / day-30 / day-90 checkpoints.",
  notes=("SIE is the sharpest hook: it's GREENFIELD and still in methodology — the expert can "
         "author its scientific spine rather than inherit it. That's the 'join now vs later' lever "
         "made concrete: in 90 days the foundational calls are made, and you want their fingerprints "
         "on them. Let them pick where to start."))

# ============================== SECTION 10 — CLOSING
S(kind="section", num="10", title="Closing Vision",
  sub="What Bari becomes if this works — and why the moment is now.",
  notes=("Land the plane. Return to ambition, but earned by everything shown. End on 'why now'."))

S(kind="content", kicker="Closing · what Bari can become",
  title="If this works",
  bullets=[
    ("A trusted consumer platform.", "The default place a shopper checks before they buy — because it's honest, not loud."),
    ("A nutrition-intelligence company.", "An engine others build on, across categories and shelves."),
    ("A professional decision-support platform.", "Where dietitians and clinicians get defensible, evidence-linked answers."),
    ("A retailer-intelligence layer.", "Category truth that retailers and manufacturers pay to understand."),
    ("International.", "The same engine, proven portable, on the US shelf and beyond."),
  ],
  visual="The 5-layer 'if this works' pyramid, each layer tagged to the section that proves it's reachable.",
  notes=("Every layer here was made credible earlier in the deck — point back. This isn't a wish "
         "list; it's the natural extension of a system that already exists and already withholds "
         "what it can't defend."))

S(kind="content", kicker="Closing · the ask",
  title="Why joining now beats joining later",
  lead="The foundational scientific decisions are being made in the next two quarters. After that, they harden.",
  bullets=[
    ("The methodology is still wet cement.", "SIE, dairy invariants, evidence grading — author them now or inherit them later."),
    ("Earliest credibility is worth the most.", "The named scientist who shaped the engine, not the one who reviewed it once it was set."),
    ("Leverage, not labour.", "Your judgement gets encoded, versioned, and applied at scale by a system that already runs."),
    ("The trust window is open now.", "The authority vacuum is real today; whoever fills it credibly first, wins."),
  ],
  visual="A 'cement setting' metaphor: foundational decisions hardening over a timeline, with a 'you are here' marker at the soft end.",
  notes=("Close on partnership, not employment: 'We've built the hard engineering and the "
         "governance. What we don't have — and won't fake — is a senior scientific conscience "
         "shaping it from the inside. That seat is open now, and it won't be this open again. "
         "I'd like it to be yours.' Then stop talking and let them respond."))

# ============================== APPENDIX
S(kind="section", num="A", title="Appendix — Deep Architecture",
  sub="For the technically curious: how BSIP0/1/2, the router, and governance actually fit.",
  notes=("Only walk these if asked. They exist to reward a deep-diving scientist and to prove "
         "there's substance under every claim."))

S(kind="content", kicker="Appendix · A1",
  title="BSIP0 — acquisition & label parsing",
  bullets=[
    ("Real-retailer ingestion.", "Shufersal scraping with structured nutrition-list parsing."),
    ("Raw-source persistence.", "Every scrape stores its raw nutrition source for offline replay and audit (TASK-151)."),
    ("Hardened parsing.", "Shared bsip0_nutrition layer; the EV-029 fat-overwrite trap fixed centrally with a permanent QA guard (COV-006)."),
    ("Output.", "Normalised per-product records: nutrition panel, ingredients, claims, provenance."),
  ],
  visual="BSIP0 data-flow: retailer → parser → normalised record, with the raw-source store branching off.",
  notes=("The audit-replay capability is the appendix highlight: we can re-derive any score from "
         "the exact raw label that produced it. Few teams can."))

S(kind="content", kicker="Appendix · A2",
  title="BSIP1 — consolidation & trust layer",
  bullets=[
    ("Cross-retailer consolidation.", "One canonical product from many noisy observations."),
    ("Trust scoring.", "Observation-quality + canonical trust weights (core/trust.py) — confidence starts here, not at scoring."),
    ("Semantic enrichment.", "10 derived fields (matrix, additive identity, fermentation signals…) that the scorer later reads."),
    ("Explicitly NOT scoring.", "BSIP1 prepares truth; BSIP2 judges it. The separation is deliberate and auditable."),
  ],
  visual="BSIP1 funnel: many observations → trust-weighted canonical product → enriched feature set.",
  notes=("Reinforce the separation-of-concerns theme: consolidation and judgement are different "
         "jobs done by different layers, so the expert can audit judgement in isolation."))

S(kind="content", kicker="Appendix · A3",
  title="BSIP2 — the scoring engine",
  columns=[
    {"h":"Four interpretive layers", "items":[
        ("Structural.", "What the food fundamentally is."),
        ("Nutritional.", "Macro/micro reality, category-relative."),
        ("Metabolic.", "How the body is likely to meet it."),
        ("Engineering.", "How processed/engineered it is — NOVA-aware, additive-identity-aware."),
    ]},
    {"h":"Universal core + archetypes", "items":[
        ("3-stage router.", "anchor → signal → resolution selects the archetype."),
        ("Archetype interpretation.", "Category-specific lenses over one shared core. Engine v0.4.0."),
        ("Confidence + null.", "Insufficient evidence → no score, by design."),
        ("Synthesis layer.", "GSS/fiber/fermentation nuance integrated post-score, governed."),
    ]},
  ],
  visual="Engine schematic: router on the left selecting an archetype, four interpretive layers stacked, confidence gate on the right.",
  notes=("If the scientist wants to go deep, this is where category calibration lives — and where "
         "their judgement would most directly write into the system. Offer to walk a single "
         "product end-to-end live if they're interested."))

S(kind="content", kicker="Appendix · A4",
  title="Governance, evidence & operating model",
  columns=[
    {"h":"Governance & evidence", "items":[
        ("Comparison constitution.", "6 articles: philosophy, eligibility, ranking, explanation, distortions, launch."),
        ("Evidence registry.", "Findings + guardrails + deferred domains; every rule traceable."),
        ("Exception registry.", "Every deliberate deviation is logged and approved before shipping."),
        ("Frozen invariants.", "Dated rulings (e.g. milk run_004) that cannot silently drift."),
    ]},
    {"h":"Operating model", "items":[
        ("Command Center.", "Authoritative task registry → derived dashboard."),
        ("Agent OS.", "Specialist agents + a verification gate before any close."),
        ("Sign-off gates.", "Roadmap-impact items can't close without required co-sign."),
        ("External evidence layer.", "10 authoritative read-only clients feed reasoning."),
    ]},
  ],
  visual="Two-panel: governance stack on the left, operating-model loop on the right, joined by the registry.",
  notes=("Close the appendix on governance — it's the least glamorous and most decisive asset, and "
         "the one the expert will personally preside over. End: 'this is the chair we're offering "
         "you.'"))

# ---------------------------------------------------------------- visual map
IMG_MAP = {
  "Nutrition information is structurally broken": ("problem_missing_layer.png","right"),
  "Trust is deteriorating — that's the opening":  ("problem_trust.png","right"),
  "What Bari is":                                  ("engine_surface.png","right"),
  "Three audiences + a marketplace, one engine":  ("three_layers_band.png","band"),
  "A credible five-year path":                    ("roadmap_band.png","band"),
  "A three-stage data pipeline: BSIP0 → BSIP1 → BSIP2": ("pipeline_band.png","band"),
  "Evidence, confidence and validation are built in — not bolted on": ("rigor_stack.png","right"),
  "Why this is hard to replicate":                ("moat_flywheel.png","right"),
  "Hummus — defining the prepared/raw boundary correctly": ("hummus_boundary.png","right"),
  "Yogurts — crediting live cultures, and refusing to overclaim": ("yogurt_grades.png","right"),
  "Cereals — testing assumptions instead of inheriting them": ("cereals_fortification.png","right"),
  "Dairy & fermentation — calibration with frozen invariants": ("dairy_fermentation.png","band"),
  "What's live today":                            ("progress_grades.png","right"),
  "Six stages, each de-risking the next":         ("monetization_staircase.png","band"),
  "How nutrition expertise plugs into the operating model": ("expert_loop.png","band"),
  "Five places you'd make an immediate difference": ("swimlane_90day.png","band"),
  "If this works":                                ("closing_pyramid.png","right"),
  "Why joining now beats joining later":          ("why_now_cement.png","right"),
}
for _s in SLIDES:
    if _s.get("title") in IMG_MAP:
        _s["img"], _s["img_mode"] = IMG_MAP[_s["title"]]

# ---- review cuts (owner pass, 2026-06-03): kill these slides. Kept in source
#      (skip=True) so they are one edit to restore.
KILL_TITLES = {
    "Cereals — testing assumptions instead of inheriting them",      # old 21
    "Dairy & fermentation — calibration with frozen invariants",     # old 22
    "Six compounding moats",                                         # old 27
    "What we need from nutrition expertise",                         # old 31
    "What is NOT expected",                                          # old 32
    "How nutrition expertise plugs into the operating model",        # old 33
    "Five places you'd make an immediate difference",                # old 35
    "If this works",                                                 # old 37
    "Why joining now beats joining later",                           # old 38
    "BSIP0 — acquisition & label parsing",                           # old 40
    "BSIP1 — consolidation & trust layer",                           # old 41
    "BSIP2 — the scoring engine",                                    # old 42
    "Governance, evidence & operating model",                        # old 43
}
KILL_SECTIONS = {"6", "8", "9", "10", "A"}   # Why-Win, Role, Opportunities, Closing, Appendix
for _s in SLIDES:
    if _s.get("title") in KILL_TITLES:
        _s["skip"] = True
    if _s["kind"] == "section" and _s.get("num") in KILL_SECTIONS:
        _s["skip"] = True
    # renumber the surviving Monetization section 7 -> 6 (sections now read 1..6)
    if _s["kind"] == "section" and _s.get("num") == "7":
        _s["num"] = "6"

def add_master_logo(prs):
    """Add the Bari mark to the slide master so it's part of the template.
    MasterShapes has no high-level add_picture, so we use the same low-level
    spTree.add_pic call python-pptx uses internally."""
    master = prs.slide_masters[0]
    shapes = master.shapes
    image_part, rId = master.part.get_or_add_image_part(LOGO_LIGHT)
    w = Inches(1.7); h = Emu(int(w * _aspect(LOGO_LIGHT)))
    x = Inches(11.18); y = Inches(0.30)
    sid = shapes._next_shape_id
    pic = shapes._spTree.add_pic(sid, "Bari Logo (master)", "Bari", rId, x, y, w, h)
    return pic

# ================================================================ BUILD
def build_pptx(path):
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    add_master_logo(prs)
    n = 0
    for s in SLIDES:
        if s.get("skip"): continue
        n += 1
        k = s["kind"]
        if k == "cover":   sl = cover_slide(prs, s)
        elif k == "section": sl = section_slide(prs, s)
        else: sl = content_slide(prs, s)
        add_pageno(sl, n, dark=(k in ("cover", "section")))
    prs.save(path)
    return len(prs.slides._sldIdLst)

def build_md(path):
    lines = ["# Bari — Nutrition Expert Partnership Deck",
             "_Slide-by-slide spec. Generated from build_deck.py (single source of truth)._\n",
             "**Grounding:** all figures trace to the live frontend corpus, the TASK registry "
             "(TASK-170 external clients, TASK-171 SIE), the evidence registry (EV-024, EV-029, "
             "COV-006) and frozen invariants (milk run_004). No product/nutrition data invented.\n",
             "---\n"]
    n = 0
    for s in SLIDES:
        if s.get("skip"): continue
        n += 1
        if s["kind"] == "cover":
            lines.append(f"## Slide {n} — COVER · {s['title']}")
            lines.append(f"- **Subtitle:** {s['subtitle']}")
            lines.append(f"- **Footer:** {s['footer']}")
            lines.append(f"- **Speaker notes:** {s['notes']}\n")
        elif s["kind"] == "section":
            lines.append(f"## Slide {n} — SECTION {s['num']} · {s['title']}")
            lines.append(f"- **Objective / sub:** {s['sub']}")
            lines.append(f"- **Speaker notes:** {s['notes']}\n")
        else:
            lines.append(f"## Slide {n} — {s['title']}")
            if s.get("kicker"): lines.append(f"- **Section:** {s['kicker']}")
            if s.get("lead"):   lines.append(f"- **Lead:** {s['lead']}")
            # bullets
            def fmt(b):
                return f"**{b[0]}** {b[1]}" if isinstance(b, tuple) else b
            if s.get("bullets"):
                lines.append("- **Key messages:**")
                for b in s["bullets"]:
                    lines.append(f"    - {fmt(b)}")
            if s.get("columns"):
                lines.append("- **Key messages (columns):**")
                for c in s["columns"]:
                    if c["h"]: lines.append(f"    - _{c['h']}_")
                    for b in c["items"]:
                        lines.append(f"        - {fmt(b)}")
            if s.get("stats"):
                lines.append("- **Headline stats:** " +
                             " · ".join(f"**{x['big']}** {x['cap']}" for x in s["stats"]))
            if s.get("visual"): lines.append(f"- **Suggested visual:** {s['visual']}")
            lines.append(f"- **Speaker notes:** {s['notes']}\n")
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    return n

if __name__ == "__main__":
    pptx_path = os.path.join(OUT_DIR, "Bari_Nutrition_Partnership.pptx")
    md_path   = os.path.join(OUT_DIR, "Bari_Nutrition_Partnership_spec.md")
    n_slides = build_pptx(pptx_path)
    n_md = build_md(md_path)
    print(f"PPTX slides: {n_slides}  ->  {pptx_path}")
    print(f"MD slides:   {n_md}  ->  {md_path}")
